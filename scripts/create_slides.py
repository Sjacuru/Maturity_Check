"""
Create the PPP Maturity Check slide deck (PPTX).

Bilingual PT | EN headings throughout.
16 slides organized as agreed during the interview phase.

Output: docs/PPP_Maturity_Check_Slides.pptx
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Theme colours
# ---------------------------------------------------------------------------
BLUE_DARK  = RGBColor(0x1A, 0x3A, 0x5C)   # slide background / title bar
BLUE_MID   = RGBColor(0x2E, 0x6D, 0xA4)   # section accents
BLUE_LIGHT = RGBColor(0xD0, 0xE4, 0xF7)   # table header background
GOLD       = RGBColor(0xC9, 0xA0, 0x2B)   # highlight / accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1C, 0x1C, 0x1C)
GREY_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
GREY_MID   = RGBColor(0xCC, 0xCC, 0xCC)
GREEN_OK   = RGBColor(0x2E, 0x7D, 0x32)
ORANGE     = RGBColor(0xE6, 0x5C, 0x00)
RED_LOW    = RGBColor(0xC6, 0x28, 0x28)

# Slide dimensions (widescreen 16:9)
W = Inches(13.33)
H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def _add_rect(slide, left, top, width, height, fill: RGBColor | None = None,
              line: RGBColor | None = None) -> object:
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height, text: str,
                  bold=False, italic=False, size=18,
                  color: RGBColor = DARK_TEXT,
                  align=PP_ALIGN.LEFT, wrap=True) -> object:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return txBox


def _title_bar(slide, title_pt: str, title_en: str,
               bar_color: RGBColor = BLUE_DARK) -> None:
    """Draw a full-width title bar at the top with bilingual title."""
    _add_rect(slide, 0, 0, W, Inches(1.1), fill=bar_color)
    _add_text_box(
        slide, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.45),
        title_pt, bold=True, size=24, color=WHITE
    )
    _add_text_box(
        slide, Inches(0.4), Inches(0.58), Inches(12.5), Inches(0.38),
        title_en, bold=False, italic=True, size=14, color=BLUE_LIGHT
    )


def _bullet_block(slide, left, top, width, height,
                  items: list[str], size=14, color=DARK_TEXT,
                  bullet_char="▸  ") -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = bullet_char + item
        run.font.size = Pt(size)
        run.font.color.rgb = color


def _section_label(slide, label: str) -> None:
    """Small label in the bottom-right corner."""
    _add_text_box(
        slide, Inches(10.5), Inches(7.1), Inches(2.5), Inches(0.3),
        label, size=9, color=GREY_MID, align=PP_ALIGN.RIGHT
    )


def _divider(slide, top, color=BLUE_MID, thickness_pt=1.5):
    from pptx.util import Pt as _Pt
    ln = slide.shapes.add_shape(1, Inches(0.3), top, Inches(12.73), Inches(0.02))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()


def _add_table(slide, left, top, width, rows, cols,
               col_widths=None, row_height=Inches(0.42)) -> object:
    table = slide.shapes.add_table(rows, cols, left, top, width,
                                   row_height * rows).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    return table


def _cell(table, row, col, text, bold=False, size=12,
          bg: RGBColor | None = None, color: RGBColor = DARK_TEXT,
          align=PP_ALIGN.LEFT) -> None:
    cell = table.cell(row, col)
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color
    if bg:
        from pptx.oxml.ns import qn
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", str(bg))


# ---------------------------------------------------------------------------
# Individual slides
# ---------------------------------------------------------------------------

def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full background
    _add_rect(slide, 0, 0, W, H, fill=BLUE_DARK)

    # Gold accent strip
    _add_rect(slide, 0, Inches(3.2), W, Inches(0.06), fill=GOLD)

    # Main title
    _add_text_box(
        slide, Inches(1.2), Inches(1.6), Inches(11), Inches(1.2),
        "PPP Maturity Check", bold=True, size=44, color=WHITE,
        align=PP_ALIGN.CENTER
    )
    _add_text_box(
        slide, Inches(1.2), Inches(2.6), Inches(11), Inches(0.6),
        "Sistema de Suporte à Decisão para Avaliação de Maturidade de PPPs",
        bold=False, size=18, color=BLUE_LIGHT, align=PP_ALIGN.CENTER
    )
    _add_text_box(
        slide, Inches(1.2), Inches(3.1), Inches(11), Inches(0.45),
        "Decision-Support System for PPP Maturity Assessment",
        bold=False, italic=True, size=14, color=BLUE_LIGHT, align=PP_ALIGN.CENTER
    )

    # Bottom metadata
    _add_text_box(
        slide, Inches(1.2), Inches(5.5), Inches(11), Inches(0.4),
        "Salim Jacuru  ·  sjacuru@gmail.com  ·  Junho 2026",
        size=12, color=GREY_MID, align=PP_ALIGN.CENTER
    )
    _add_text_box(
        slide, Inches(1.2), Inches(5.9), Inches(11), Inches(0.35),
        "Framework IPMP (TCU 2026)  ·  46 Ações  ·  5 Dimensões  ·  Município do Rio de Janeiro",
        size=11, color=GREY_MID, align=PP_ALIGN.CENTER
    )


def slide_agenda(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide, "Estrutura da Apresentação", "Presentation Structure")

    cols = [
        ("1  Problema e Contexto", "Problem and Context", BLUE_DARK),
        ("2  A Solução", "The Solution", BLUE_MID),
        ("3  Arquitetura", "Architecture", BLUE_DARK),
        ("4  Módulos", "Modules", BLUE_MID),
        ("5  Fluxo do Auditor", "Auditor Workflow", BLUE_DARK),
        ("6  Estado Atual", "Current State", BLUE_MID),
        ("7  Validação (Phase A)", "Validation (Phase A)", BLUE_DARK),
        ("8  Roteiro", "Roadmap", BLUE_MID),
    ]
    col_w = Inches(3.1)
    col_h = Inches(1.3)
    margin_left = Inches(0.4)
    gap = Inches(0.18)
    top = Inches(1.3)

    for i, (pt, en, color) in enumerate(cols):
        row, col = divmod(i, 4)
        left = margin_left + col * (col_w + gap)
        t = top + row * (col_h + gap)
        _add_rect(slide, left, t, col_w, col_h, fill=color)
        _add_text_box(slide, left + Inches(0.15), t + Inches(0.15),
                      col_w - Inches(0.3), Inches(0.5),
                      pt, bold=True, size=14, color=WHITE)
        _add_text_box(slide, left + Inches(0.15), t + Inches(0.65),
                      col_w - Inches(0.3), Inches(0.45),
                      en, italic=True, size=11, color=BLUE_LIGHT)

    _section_label(slide, "Estrutura | Structure")


def slide_problema(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "O Problema: avaliação de maturidade de projetos PPP",
               "The Problem: maturity assessment of PPP projects")

    # Three columns: Volume / Subjectividade / Reprodutibilidade
    boxes = [
        ("Volume", "Dezenas de documentos, centenas de páginas por processo",
         "Dozens of documents, hundreds of pages per process"),
        ("Subjetividade", "Avaliadores diferentes → scores diferentes sem rastreabilidade",
         "Different evaluators → different scores without traceability"),
        ("Reprodutibilidade", "Requisito acadêmico: mesma entrada deve gerar mesmo resultado",
         "Academic requirement: same input must produce same result"),
    ]
    col_w = Inches(3.9)
    gap = Inches(0.28)
    left0 = Inches(0.4)
    top = Inches(1.35)
    for i, (title, pt, en) in enumerate(boxes):
        left = left0 + i * (col_w + gap)
        _add_rect(slide, left, top, col_w, Inches(4.8), fill=GREY_LIGHT)
        _add_rect(slide, left, top, col_w, Inches(0.55), fill=BLUE_DARK)
        _add_text_box(slide, left + Inches(0.12), top + Inches(0.08),
                      col_w - Inches(0.2), Inches(0.4),
                      title, bold=True, size=16, color=WHITE)
        _add_text_box(slide, left + Inches(0.15), top + Inches(0.7),
                      col_w - Inches(0.3), Inches(1.8),
                      pt, size=13, color=DARK_TEXT)
        _divider(slide, top + Inches(2.6), color=GREY_MID)
        _add_text_box(slide, left + Inches(0.15), top + Inches(2.7),
                      col_w - Inches(0.3), Inches(1.8),
                      en, size=11, italic=True, color=BLUE_MID)

    _section_label(slide, "1  Problema | Problem")


def slide_ipmp(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "O Framework IPMP — 46 Ações, 5 Dimensões",
               "The IPMP Framework — 46 Actions, 5 Dimensions")

    # Dimensions table
    dims = [
        ("1", "Estratégica", "Strategic"),
        ("2", "Técnica", "Technical"),
        ("3", "Financeira", "Financial"),
        ("4", "Ambiental e Social", "Environmental and Social"),
        ("5", "Jurídica e Regulatória", "Legal and Regulatory"),
    ]
    t = _add_table(slide, Inches(0.4), Inches(1.3), Inches(7.5),
                   rows=6, cols=3,
                   col_widths=[Inches(0.6), Inches(3.3), Inches(3.6)],
                   row_height=Inches(0.46))
    _cell(t, 0, 0, "#", bold=True, bg=BLUE_DARK, color=WHITE, align=PP_ALIGN.CENTER)
    _cell(t, 0, 1, "Dimensão (PT)", bold=True, bg=BLUE_DARK, color=WHITE)
    _cell(t, 0, 2, "Dimension (EN)", bold=True, bg=BLUE_DARK, color=WHITE)
    for i, (num, pt, en) in enumerate(dims, 1):
        bg = BLUE_LIGHT if i % 2 == 0 else WHITE
        _cell(t, i, 0, num, bold=True, bg=bg, align=PP_ALIGN.CENTER)
        _cell(t, i, 1, pt, bg=bg)
        _cell(t, i, 2, en, bg=bg)

    # Score box
    _add_rect(slide, Inches(8.2), Inches(1.3), Inches(4.7), Inches(2.7),
              fill=BLUE_DARK)
    _add_text_box(slide, Inches(8.4), Inches(1.4), Inches(4.3), Inches(0.45),
                  "Pontuação | Score", bold=True, size=14, color=WHITE)
    for val, label_pt, label_en, col in [
        ("0", "Não Atendido", "Not Met", RED_LOW),
        ("1", "Parcialmente Atendido", "Partially Met", ORANGE),
        ("3", "Atendido", "Met", GREEN_OK),
    ]:
        pass
    score_rows = [
        ("0", "Não Atendido", "Not Met"),
        ("1", "Parcialmente Atendido", "Partially Met"),
        ("3", "Atendido", "Met"),
    ]
    colors_score = [RED_LOW, ORANGE, GREEN_OK]
    for i, ((val, pt, en), col) in enumerate(zip(score_rows, colors_score)):
        top_s = Inches(1.95) + i * Inches(0.62)
        _add_rect(slide, Inches(8.3), top_s, Inches(0.55), Inches(0.5), fill=col)
        _add_text_box(slide, Inches(8.3), top_s + Inches(0.08),
                      Inches(0.55), Inches(0.35),
                      val, bold=True, size=18, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text_box(slide, Inches(8.95), top_s + Inches(0.05),
                      Inches(3.8), Inches(0.25),
                      pt, size=12, color=WHITE)
        _add_text_box(slide, Inches(8.95), top_s + Inches(0.27),
                      Inches(3.8), Inches(0.22),
                      en, size=10, italic=True, color=BLUE_LIGHT)

    # Max score
    _add_rect(slide, Inches(8.2), Inches(4.1), Inches(4.7), Inches(0.7),
              fill=GOLD)
    _add_text_box(slide, Inches(8.3), Inches(4.15), Inches(4.5), Inches(0.6),
                  "Pontuação máxima: 138 pts  (46 × 3)",
                  bold=True, size=14, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    # Scope note
    _add_rect(slide, Inches(0.4), Inches(4.6), Inches(12.5), Inches(0.7),
              fill=BLUE_LIGHT)
    _add_text_box(slide, Inches(0.6), Inches(4.65), Inches(12), Inches(0.6),
                  "Escopo Phase 1: Ação 1 apenas — \"Descreva o projeto, seu contexto e os objetivos estratégicos\"  |  "
                  "Phase 1 scope: Action 1 only — \"Describe the project, its context, and strategic objectives\"",
                  size=12, color=BLUE_DARK)

    _section_label(slide, "1  Problema | Problem")


def slide_solucao(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "A Solução: suporte à decisão, não substituição do auditor",
               "The Solution: decision support, not auditor replacement")

    # Big concept
    _add_rect(slide, Inches(0.4), Inches(1.25), Inches(12.5), Inches(1.3),
              fill=BLUE_DARK)
    _add_text_box(slide, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5),
                  "O sistema recupera evidências, raciocina e propõe uma pontuação.",
                  bold=True, size=18, color=WHITE)
    _add_text_box(slide, Inches(0.6), Inches(1.82), Inches(12), Inches(0.4),
                  "The system retrieves evidence, reasons, and proposes a score.  The auditor validates.",
                  italic=True, size=13, color=BLUE_LIGHT)

    # Four steps flow
    steps = [
        ("1", "Recuperação", "Retrieval", "Localiza trechos\nrelevantes\nno processo",
         "Locates relevant\npassages\nin the process"),
        ("2", "Avaliação", "Evaluation", "LLM raciocina\nsobre as evidências\ncom critérios IPMP",
         "LLM reasons\nabout evidence\nusing IPMP criteria"),
        ("3", "Apresentação", "Presentation", "7 elementos do\npacote de evidências\npara o auditor",
         "7-element evidence\npackage for\nthe auditor"),
        ("4", "Validação", "Validation", "Auditor aceita\nou substitui\ncom justificativa",
         "Auditor accepts\nor overrides\nwith justification"),
    ]
    box_w = Inches(2.8)
    gap = Inches(0.55)
    left0 = Inches(0.55)
    top = Inches(2.75)
    for i, (num, pt, en, desc_pt, desc_en) in enumerate(steps):
        left = left0 + i * (box_w + gap)
        _add_rect(slide, left, top, box_w, Inches(3.8), fill=GREY_LIGHT)
        _add_rect(slide, left, top, box_w, Inches(0.58), fill=BLUE_MID)
        _add_text_box(slide, left, top + Inches(0.1),
                      Inches(0.6), Inches(0.4),
                      num, bold=True, size=18, color=GOLD, align=PP_ALIGN.CENTER)
        _add_text_box(slide, left + Inches(0.55), top + Inches(0.1),
                      box_w - Inches(0.65), Inches(0.28),
                      pt, bold=True, size=14, color=WHITE)
        _add_text_box(slide, left + Inches(0.55), top + Inches(0.36),
                      box_w - Inches(0.65), Inches(0.22),
                      en, italic=True, size=10, color=BLUE_LIGHT)
        _add_text_box(slide, left + Inches(0.15), top + Inches(0.75),
                      box_w - Inches(0.3), Inches(1.5),
                      desc_pt, size=12, color=DARK_TEXT)
        _divider(slide, top + Inches(2.3), color=GREY_MID)
        _add_text_box(slide, left + Inches(0.15), top + Inches(2.4),
                      box_w - Inches(0.3), Inches(1.1),
                      desc_en, size=10, italic=True, color=BLUE_MID)

        # Arrow between steps
        if i < 3:
            arr_left = left + box_w + Inches(0.12)
            _add_text_box(slide, arr_left, top + Inches(1.5),
                          Inches(0.35), Inches(0.45),
                          "→", bold=True, size=22, color=BLUE_MID,
                          align=PP_ALIGN.CENTER)

    _section_label(slide, "2  Solução | Solution")


def slide_fluxo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "Fluxo de Informação Ponta-a-Ponta",
               "End-to-End Information Flow")

    stages = [
        ("PDF do processo\n(EVTEA, Relatório...)",
         "Process PDF\n(EVTEA, Report...)", BLUE_DARK),
        ("Chunks estruturados\n(texto + metadados)",
         "Structured chunks\n(text + metadata)", BLUE_MID),
        ("Evidências recuperadas\n(com proveniência)",
         "Retrieved evidence\n(with provenance)", BLUE_DARK),
        ("Raciocínio + score\nproposto",
         "Reasoning + proposed\nscore", BLUE_MID),
        ("Decisão do auditor\n(final + justificativa)",
         "Auditor decision\n(final + justification)", BLUE_DARK),
    ]
    modules = [
        "Módulo 2\nExtraction",
        "Módulo 3\nRetrieval",
        "Módulo 4\nEvaluation",
        "Módulo 5\nAssessment",
    ]
    box_w = Inches(2.25)
    box_h = Inches(1.3)
    gap = Inches(0.18)
    left0 = Inches(0.3)
    top_box = Inches(1.6)
    top_arr = Inches(2.0)
    top_mod = Inches(3.05)

    for i, (pt, en, color) in enumerate(stages):
        left = left0 + i * (box_w + gap)
        _add_rect(slide, left, top_box, box_w, box_h, fill=color)
        _add_text_box(slide, left + Inches(0.1), top_box + Inches(0.1),
                      box_w - Inches(0.2), Inches(0.58),
                      pt, bold=True, size=12, color=WHITE)
        _add_text_box(slide, left + Inches(0.1), top_box + Inches(0.68),
                      box_w - Inches(0.2), Inches(0.5),
                      en, italic=True, size=10, color=BLUE_LIGHT)

        if i < 4:
            arr_left = left + box_w + Inches(0.02)
            _add_text_box(slide, arr_left, top_arr,
                          Inches(0.17), Inches(0.5),
                          "→", bold=True, size=20, color=BLUE_MID,
                          align=PP_ALIGN.CENTER)
            mod_left = left + box_w + Inches(0.0)
            _add_rect(slide, mod_left, top_mod,
                      Inches(0.17), Inches(0.72), fill=GOLD)

    for i, mod_text in enumerate(modules):
        left = left0 + (i + 0.5) * (box_w + gap) + box_w - Inches(0.7)
        _add_text_box(slide, left, top_mod, Inches(1.5), Inches(0.7),
                      mod_text, size=9, color=BLUE_DARK, bold=True)

    # Principles strip
    _add_rect(slide, Inches(0.3), Inches(4.3), Inches(12.7), Inches(2.8),
              fill=GREY_LIGHT)
    _add_text_box(slide, Inches(0.5), Inches(4.4), Inches(12), Inches(0.4),
                  "Princípios de Projeto | Design Principles",
                  bold=True, size=14, color=BLUE_DARK)
    principles = [
        ("Reprodutibilidade | Reproducibility",
         "BM25 determinístico + temperatura=0 → mesma entrada, mesmo score\n"
         "Deterministic BM25 + temperature=0 → same input, same score"),
        ("Rastreabilidade | Traceability",
         "Cada score carrega sua cadeia de evidências completa\n"
         "Every score carries its complete evidence chain"),
        ("Humano no loop | Human in the loop",
         "A pontuação final é sempre do auditor, nunca do sistema\n"
         "The final score always belongs to the auditor, never the system"),
    ]
    for i, (title, desc) in enumerate(principles):
        left_p = Inches(0.5) + i * Inches(4.2)
        _add_text_box(slide, left_p, Inches(4.85), Inches(4.0), Inches(0.3),
                      title, bold=True, size=11, color=BLUE_MID)
        _add_text_box(slide, left_p, Inches(5.18), Inches(4.0), Inches(1.6),
                      desc, size=10, color=DARK_TEXT)

    _section_label(slide, "3  Arquitetura | Architecture")


def slide_modulos(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "Mapa de Módulos e Responsabilidades",
               "Module Responsibility Map")

    rows_data = [
        ("1 — Ingestion", "Referência", "Reference",
         "Arquivos IPMP + Rio Manual", "IPMP + Rio Manual files",
         "Store de critérios e metadados", "Criteria and search metadata store"),
        ("2 — Extraction", "Extração", "Extraction",
         "PDF do processo", "Process PDF",
         "Lista de Chunks estruturados", "List of structured Chunks"),
        ("3 — Retrieval", "Recuperação", "Retrieval",
         "Chunks indexados + critérios", "Indexed chunks + criteria",
         "RetrievedChunks com proveniência", "RetrievedChunks with provenance"),
        ("4 — Evaluation", "Avaliação", "Evaluation",
         "RetrievedChunks + critérios IPMP", "RetrievedChunks + IPMP criteria",
         "EvaluationResult + raciocínio", "EvaluationResult + reasoning"),
        ("5 — Assessment", "Orquestração + API", "Orchestration + API",
         "PDFs do caso", "Case PDFs",
         "Resultados persistidos + REST API", "Persisted results + REST API"),
        ("6 — Vector Fallback", "Recuperação (fallback)", "Retrieval (fallback)",
         "Chunks sem resultado léxico", "Chunks with no lexical result",
         "Chunks por similaridade semântica", "Chunks by semantic similarity"),
        ("7 — Frontend", "Interface do auditor", "Auditor interface",
         "API REST", "REST API",
         "Painel Vue.js com evidências", "Vue.js panel with evidence"),
    ]

    t = _add_table(slide, Inches(0.3), Inches(1.25), Inches(12.7),
                   rows=len(rows_data) + 1, cols=4,
                   col_widths=[Inches(2.2), Inches(2.1), Inches(3.6), Inches(4.8)],
                   row_height=Inches(0.72))
    headers = ["Módulo | Module", "Estágio | Stage",
               "Entrada | Input", "Saída | Output"]
    for c, h in enumerate(headers):
        _cell(t, 0, c, h, bold=True, bg=BLUE_DARK, color=WHITE, size=11)

    for r, (mod, stg_pt, stg_en, inp_pt, inp_en, out_pt, out_en) in enumerate(rows_data, 1):
        bg = BLUE_LIGHT if r % 2 == 0 else WHITE
        _cell(t, r, 0, mod, bold=True, bg=bg, size=10)
        _cell(t, r, 1, f"{stg_pt}\n{stg_en}", bg=bg, size=10)
        _cell(t, r, 2, f"{inp_pt}\n{inp_en}", bg=bg, size=10)
        _cell(t, r, 3, f"{out_pt}\n{out_en}", bg=bg, size=10)

    _section_label(slide, "4  Módulos | Modules")


def slide_cascata(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "Módulo 3 — Retrieval: Cascata de Recuperação",
               "Module 3 — Retrieval: Retrieval Cascade")

    steps = [
        ("A", "Filename Match",
         "Compara nomes de arquivo com documentos do Rio Manual\n"
         "Matches file names against Rio Manual document names",
         "Mais confiável | Most reliable"),
        ("B", "Variant Match",
         "Busca variantes do nome do documento nos primeiros chunks\n"
         "Searches document name variants in first chunks",
         ""),
        ("C", "BM25 Corpus-Wide",
         "FTS5 OR queries: termos IPMP + hints curados Rio Manual\n"
         "FTS5 OR queries: IPMP terms + curated Rio Manual hints",
         "⚠  Fix A1 aplicado | applied"),
        ("D", "Regex Patterns",
         "Padrões legais específicos (ex: Lei Complementar n.º 105/2009)\n"
         "Specific legal patterns (e.g., Complementary Law No. 105/2009)",
         "Aditivo ao BM25 | Additive to BM25"),
        ("E", "Vector Fallback",
         "LanceDB + all-MiniLM-L6-v2 — apenas se A–D = ∅\n"
         "LanceDB + all-MiniLM-L6-v2 — only if A–D = ∅",
         "Menos confiável | Least reliable"),
    ]
    step_h = Inches(0.98)
    top0 = Inches(1.3)
    left_box = Inches(0.4)
    for i, (letter, title, desc, note) in enumerate(steps):
        top = top0 + i * (step_h + Inches(0.06))
        color = BLUE_DARK if i % 2 == 0 else BLUE_MID
        _add_rect(slide, left_box, top, Inches(0.65), step_h, fill=color)
        _add_text_box(slide, left_box, top + Inches(0.2),
                      Inches(0.65), Inches(0.55),
                      letter, bold=True, size=22, color=WHITE,
                      align=PP_ALIGN.CENTER)
        _add_rect(slide, left_box + Inches(0.65), top,
                  Inches(11.6), step_h, fill=GREY_LIGHT)
        _add_text_box(slide, left_box + Inches(0.8), top + Inches(0.07),
                      Inches(8.5), Inches(0.3),
                      title, bold=True, size=13, color=BLUE_DARK)
        _add_text_box(slide, left_box + Inches(0.8), top + Inches(0.38),
                      Inches(8.5), Inches(0.55),
                      desc, size=10, color=DARK_TEXT)
        if note:
            note_color = ORANGE if "Fix" in note else GREY_MID
            _add_text_box(slide, Inches(9.8), top + Inches(0.28),
                          Inches(3.8), Inches(0.45),
                          note, size=10, italic=True,
                          color=note_color, align=PP_ALIGN.RIGHT)

        if i < 4:
            arr_left = left_box + Inches(0.2)
            next_top = top + step_h + Inches(0.01)
            _add_text_box(slide, arr_left, next_top,
                          Inches(0.3), Inches(0.1),
                          "↓", size=10, color=BLUE_MID, align=PP_ALIGN.CENTER)

    _section_label(slide, "4  Módulos | Modules")


def slide_prompt(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "Módulo 4 — Evaluation: Prompt e Raciocínio",
               "Module 4 — Evaluation: Prompt and Reasoning")

    # Left: prompt structure
    _add_rect(slide, Inches(0.3), Inches(1.25), Inches(6.0), Inches(5.9),
              fill=GREY_LIGHT)
    _add_rect(slide, Inches(0.3), Inches(1.25), Inches(6.0), Inches(0.42),
              fill=BLUE_DARK)
    _add_text_box(slide, Inches(0.45), Inches(1.3),
                  Inches(5.7), Inches(0.35),
                  "Prompt de Sistema | System Prompt", bold=True, size=13, color=WHITE)

    sp_items = [
        "Papel: avaliador especialista IPMP",
        "Critérios detalhados da Ação (produtos 1a–1d)",
        "Rubric: exemplos pontuados (Atendido / Parcialmente / Não Atendido)",
        "Instrução de formato: raciocínio livre + bloco sentinela obrigatório",
        "─────────────────────────────────────────",
        "Role: expert IPMP evaluator",
        "Detailed action criteria (products 1a–1d)",
        "Rubric: scored examples (Met / Partial / Not Met)",
        "Format: free reasoning + mandatory sentinel block",
    ]
    _bullet_block(slide, Inches(0.45), Inches(1.75),
                  Inches(5.7), Inches(2.3), sp_items, size=10)

    _add_rect(slide, Inches(0.3), Inches(4.05), Inches(6.0), Inches(0.42),
              fill=BLUE_MID)
    _add_text_box(slide, Inches(0.45), Inches(4.1),
                  Inches(5.7), Inches(0.35),
                  "Prompt de Usuário | User Prompt", bold=True, size=13, color=WHITE)
    up_items = [
        "[Arquivo: Relatório de Pré-Análise.pdf | Página: 3]",
        "... texto do chunk ...",
        "[Arquivo: EVTEA.pdf | Página: 12]",
        "... texto do chunk ...",
        "(até 20.000 chars; prioridade por etapa da cascata)",
        "(up to 20,000 chars; priority by cascade step)",
    ]
    _bullet_block(slide, Inches(0.45), Inches(4.55),
                  Inches(5.7), Inches(2.3), up_items, size=10)

    # Right: parse + output
    _add_rect(slide, Inches(6.6), Inches(1.25), Inches(6.4), Inches(5.9),
              fill=GREY_LIGHT)
    _add_rect(slide, Inches(6.6), Inches(1.25), Inches(6.4), Inches(0.42),
              fill=BLUE_DARK)
    _add_text_box(slide, Inches(6.75), Inches(1.3),
                  Inches(6.1), Inches(0.35),
                  "Resposta do LLM | LLM Response", bold=True, size=13, color=WHITE)

    example = (
        "Raciocínio:\n"
        "Os Produtos Esperados 1a–1d estão evidenciados:\n"
        "1a. Necessidade → Seção 1 do documento\n"
        "1b. Contexto econômico → Seção 2\n"
        "1c. Objetivos estratégicos → Seção 3\n"
        "1d. Alinhamento com PPA/LOA/LDO → Seção 4\n\n"
        "Bloco sentinela:\n"
        "SCORE: 3\n"
        "UNCERTAINTY: no"
    )
    _add_rect(slide, Inches(6.75), Inches(1.75), Inches(6.1), Inches(2.55),
              fill=BLUE_DARK, line=BLUE_MID)
    _add_text_box(slide, Inches(6.85), Inches(1.8),
                  Inches(5.9), Inches(2.45),
                  example, size=10, color=WHITE)

    _add_text_box(slide, Inches(6.75), Inches(4.4),
                  Inches(6.0), Inches(0.35),
                  "Modo de falha seguro | Safe failure mode", bold=True,
                  size=12, color=BLUE_DARK)
    fail_items = [
        "parse_failed=True se o sentinela não for encontrado",
        "Auditor é notificado; score não é gravado",
        "parse_failed=True if sentinel not found",
        "Auditor is notified; score is not recorded",
    ]
    _bullet_block(slide, Inches(6.75), Inches(4.8),
                  Inches(6.0), Inches(1.8), fail_items, size=10)

    _section_label(slide, "4  Módulos | Modules")


def slide_auditor(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "Fluxo de Trabalho do Auditor",
               "Auditor Workflow")

    # Journey steps
    steps = [
        ("Upload", "Informa número do processo\ne faz upload dos PDFs",
         "Enters process number\nand uploads PDFs"),
        ("Indexação", "Sistema extrai, chunka\ne indexa o texto",
         "Extraction\nchunking\nindexing"),
        ("Recuperação\n+ Avaliação", "Cascata BM25/regex;\nLLM raciocina",
         "BM25/regex cascade;\nLLM reasons"),
        ("Painel\nAção 1", "7 elementos do\npacote de evidências",
         "7 elements of\nevidence package"),
        ("Decisão", "Aceitar ou substituir\ncom justificativa",
         "Accept or override\nwith justification"),
    ]
    box_w = Inches(2.2)
    gap = Inches(0.42)
    left0 = Inches(0.35)
    top_j = Inches(1.35)
    for i, (title, pt, en) in enumerate(steps):
        left = left0 + i * (box_w + gap)
        color = BLUE_DARK if i % 2 == 0 else BLUE_MID
        _add_rect(slide, left, top_j, box_w, Inches(2.0), fill=color)
        _add_text_box(slide, left + Inches(0.1), top_j + Inches(0.1),
                      box_w - Inches(0.2), Inches(0.5),
                      title, bold=True, size=13, color=WHITE)
        _add_text_box(slide, left + Inches(0.1), top_j + Inches(0.62),
                      box_w - Inches(0.2), Inches(0.7),
                      pt, size=11, color=WHITE)
        _add_text_box(slide, left + Inches(0.1), top_j + Inches(1.35),
                      box_w - Inches(0.2), Inches(0.6),
                      en, italic=True, size=9, color=BLUE_LIGHT)
        if i < 4:
            arr_l = left + box_w + Inches(0.1)
            _add_text_box(slide, arr_l, top_j + Inches(0.7),
                          Inches(0.28), Inches(0.55),
                          "→", bold=True, size=20, color=BLUE_MID,
                          align=PP_ALIGN.CENTER)

    # 7-element table
    _add_text_box(slide, Inches(0.35), Inches(3.55), Inches(6), Inches(0.35),
                  "Os 7 elementos do pacote de evidências | The 7-element evidence package",
                  bold=True, size=13, color=BLUE_DARK)

    elements = [
        ("Evidências recuperadas", "Retrieved evidence"),
        ("Proveniência da recuperação", "Retrieval provenance (cascade step + query)"),
        ("Critérios IPMP completos", "Full IPMP criteria"),
        ("Prompt completo", "Full prompt (system + user)"),
        ("Raciocínio do LLM", "LLM reasoning (per expected product)"),
        ("Flag de incerteza", "Uncertainty flag"),
        ("Pontuação proposta", "Proposed score (0 / 1 / 3)"),
    ]
    t = _add_table(slide, Inches(0.35), Inches(3.95), Inches(7.8),
                   rows=len(elements) + 1, cols=2,
                   col_widths=[Inches(3.6), Inches(4.2)],
                   row_height=Inches(0.42))
    _cell(t, 0, 0, "Elemento (PT)", bold=True, bg=BLUE_DARK, color=WHITE, size=11)
    _cell(t, 0, 1, "Element (EN)", bold=True, bg=BLUE_DARK, color=WHITE, size=11)
    for i, (pt, en) in enumerate(elements, 1):
        bg = BLUE_LIGHT if i % 2 == 0 else WHITE
        _cell(t, i, 0, pt, bg=bg, size=10)
        _cell(t, i, 1, en, bg=bg, size=10)

    # Final decision box
    _add_rect(slide, Inches(8.35), Inches(3.55), Inches(4.65), Inches(3.75),
              fill=BLUE_DARK)
    _add_text_box(slide, Inches(8.5), Inches(3.65), Inches(4.3), Inches(0.4),
                  "Decisão do Auditor | Auditor Decision",
                  bold=True, size=13, color=WHITE)
    dec_items = [
        "✓  Aceitar a pontuação proposta\n     Accept proposed score",
        "",
        "↔  Substituir com nova pontuação\n     Override with new score",
        "    + justificativa obrigatória\n     + mandatory justification",
        "",
        "Score final sempre do auditor.\nFinal score always the auditor's.",
    ]
    _bullet_block(slide, Inches(8.5), Inches(4.1),
                  Inches(4.3), Inches(3.0), dec_items,
                  size=11, color=WHITE, bullet_char="")

    _section_label(slide, "5  Auditor | Workflow")


def slide_estado_atual(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "Estado Atual da Implementação",
               "Current Implementation State")

    # Stats bar
    for i, (num, label_pt, label_en) in enumerate([
        ("7", "Módulos\nconcluídos", "Completed\nmodules"),
        ("300", "Testes\npassando", "Tests\npassing"),
        ("1", "Ação IPMP\nimplementada", "IPMP Action\nimplemented"),
        ("E2E", "Fluxo ponta-a-ponta\nfuncional", "End-to-end\npipeline functional"),
    ]):
        left = Inches(0.35) + i * Inches(3.2)
        _add_rect(slide, left, Inches(1.25), Inches(3.0), Inches(1.4), fill=BLUE_MID)
        _add_text_box(slide, left + Inches(0.1), Inches(1.3),
                      Inches(2.8), Inches(0.65),
                      num, bold=True, size=36, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text_box(slide, left + Inches(0.1), Inches(1.93),
                      Inches(2.8), Inches(0.35),
                      f"{label_pt}  |  {label_en}",
                      size=9, color=WHITE, align=PP_ALIGN.CENTER)

    # Observations
    _add_text_box(slide, Inches(0.35), Inches(2.85), Inches(6.1), Inches(0.35),
                  "O que foi observado | What was observed",
                  bold=True, size=13, color=BLUE_DARK)
    obs = [
        "Smoke test com documento M5D real: parse_failed=True (score não produzido)",
        "Causa: 57.000 chars de evidência → janela de contexto Ollama esgotada",
        "Correção: num_ctx=32768 + cap de 20k chars com prioridade de cascata",
        "Fix A1: BM25 ignorava 25 termos curados do Rio Manual → corrigido",
        "──────────────────────────────────────────────────────────────────",
        "Smoke test with real M5D document: parse_failed=True (score not produced)",
        "Cause: 57,000 chars evidence → Ollama context window exhausted",
        "Fix: num_ctx=32768 + 20k char cap with cascade priority",
        "Fix A1: BM25 was ignoring 25 curated Rio Manual terms → corrected",
    ]
    _bullet_block(slide, Inches(0.35), Inches(3.25),
                  Inches(6.1), Inches(3.8), obs, size=10)

    # Limitations
    _add_text_box(slide, Inches(6.7), Inches(2.85), Inches(6.3), Inches(0.35),
                  "Limitações atuais | Current limitations",
                  bold=True, size=13, color=BLUE_DARK)
    lims = [
        "Escopo: Ação 1 apenas (45 ações pendentes)",
        "Modelo: apenas Mistral (Ollama); sem comparação",
        "Validação: corpus sintético apenas; real-world pendente",
        "Frontend funcional, não validado com usuários reais",
        "──────────────────────────────────────────────────",
        "Scope: Action 1 only (45 actions pending)",
        "Model: Mistral (Ollama) only; no comparison",
        "Validation: synthetic corpus only; real-world pending",
        "Frontend functional, not validated with real users",
    ]
    _bullet_block(slide, Inches(6.7), Inches(3.25),
                  Inches(6.3), Inches(3.8), lims, size=10)

    _section_label(slide, "6  Estado Atual | Current State")


def slide_validacao(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "Validação Phase A — Cadeia de Raciocínio",
               "Phase A Validation — Reasoning Chain")

    _add_text_box(slide, Inches(0.35), Inches(1.25), Inches(12.5), Inches(0.35),
                  "Corpus sintético controlado com ground truth conhecido | Controlled synthetic corpus with known ground truth",
                  size=12, color=BLUE_MID, italic=True)

    # Results table
    rows_v = [
        ("A1  Retrieval", "BM25 usava apenas termos genéricos IPMP",
         "BM25 only used generic IPMP terms",
         "Fix aplicado: +25 termos Rio Manual\nFix applied: +25 Rio Manual terms",
         "Média | Medium", ORANGE),
        ("A3  Chunk Coverage", "Score3: 4/4 páginas recuperadas via BM25\n(3 via rio_hints; 1 via produto 1b)",
         "Score3: 4/4 pages via BM25\n(3 via rio_hints; 1 via product 1b)",
         "Score1: 1 chunk (linguagem vaga limita)\nScore1: 1 chunk (vague language limits retrieval)",
         "Média | Medium", ORANGE),
        ("A4  Prompt\n+ Reasoning", "Score3 → proposto=3 ✓\nparse_failed=False, uncertainty=False",
         "Score3 → proposed=3 ✓\nparse_failed=False, uncertainty=False",
         "Score1 → proposto=1 ✓\nLLM mapeou 1a→S1, 1b→S2, 1c→S3, 1d→S4",
         "Média | Medium", ORANGE),
        ("A4  Model\nComparison", "Apenas Mistral (Ollama) testado",
         "Only Mistral (Ollama) tested",
         "Groq e outros modelos não avaliados\nGroq and other models not evaluated",
         "Desconhecida | Unknown", RED_LOW),
    ]
    t = _add_table(slide, Inches(0.3), Inches(1.7), Inches(12.7),
                   rows=len(rows_v) + 1, cols=4,
                   col_widths=[Inches(1.8), Inches(3.7), Inches(4.2), Inches(3.0)],
                   row_height=Inches(1.1))
    hdrs = ["Estágio | Stage", "Achado (PT) | Finding (PT)",
            "Achado (EN) | Finding (EN)", "Confiança | Confidence"]
    for c, h in enumerate(hdrs):
        _cell(t, 0, c, h, bold=True, bg=BLUE_DARK, color=WHITE, size=10)
    for r, (stage, pt, en, extra, conf, conf_color) in enumerate(rows_v, 1):
        bg = BLUE_LIGHT if r % 2 == 0 else WHITE
        _cell(t, r, 0, stage, bold=True, bg=bg, size=10)
        _cell(t, r, 1, pt, bg=bg, size=9)
        _cell(t, r, 2, en, bg=bg, size=9)
        _cell(t, r, 3, conf, bold=True, bg=conf_color, color=WHITE,
              size=11, align=PP_ALIGN.CENTER)

    # Legend
    for label, color in [("Alta | High", GREEN_OK),
                          ("Média | Medium", ORANGE),
                          ("Baixa | Low", RED_LOW),
                          ("Desconhecida | Unknown", RED_LOW)]:
        pass
    _add_rect(slide, Inches(0.3), Inches(6.85), Inches(1.4), Inches(0.4), fill=GREEN_OK)
    _add_text_box(slide, Inches(0.35), Inches(6.9), Inches(1.3), Inches(0.3),
                  "Alta | High", size=9, color=WHITE, bold=True)
    _add_rect(slide, Inches(1.8), Inches(6.85), Inches(1.4), Inches(0.4), fill=ORANGE)
    _add_text_box(slide, Inches(1.85), Inches(6.9), Inches(1.3), Inches(0.3),
                  "Média | Medium", size=9, color=WHITE, bold=True)
    _add_rect(slide, Inches(3.3), Inches(6.85), Inches(1.4), Inches(0.4), fill=RED_LOW)
    _add_text_box(slide, Inches(3.35), Inches(6.9), Inches(1.3), Inches(0.3),
                  "Baixa / Desconhecida", size=9, color=WHITE, bold=True)

    _section_label(slide, "7  Validação | Validation")


def slide_confianca(prs: Presentation) -> None:
    """Confidence level per reasoning chain stage — the key Phase A output."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "Cadeia de Raciocínio — Confiança por Estágio",
               "Reasoning Chain — Confidence per Stage")

    stages_conf = [
        ("Estágio 1\nRetrieval", "Stage 1\nRetrieval",
         "BM25 + hints Rio Manual\nfunciona no corpus controlado",
         "BM25 + Rio Manual hints\nworks on controlled corpus",
         "MÉDIA", ORANGE),
        ("Estágio 2\nCobertura", "Stage 2\nCoverage",
         "Corpus sintético cobre 1a–1d explicitamente;\ndocumentos reais são implícitos",
         "Synthetic corpus covers 1a–1d explicitly;\nreal documents use implicit language",
         "BAIXA\npara docs reais", RED_LOW),
        ("Estágio 3\nPrompt", "Stage 3\nPrompt",
         "Raciocínio correto em Score3 e Score1;\ncasos ambíguos não testados",
         "Correct reasoning on Score3 and Score1;\nambiguous cases not tested",
         "MÉDIA", ORANGE),
        ("Estágio 4\nModelo", "Stage 4\nModel",
         "Apenas Mistral testado;\ncomparação com outros modelos pendente",
         "Only Mistral tested;\ncomparison with other models pending",
         "DESCONHECIDA", RED_LOW),
    ]
    box_w = Inches(2.9)
    gap = Inches(0.25)
    left0 = Inches(0.35)
    top_s = Inches(1.35)
    for i, (pt_s, en_s, pt_d, en_d, conf, conf_color) in enumerate(stages_conf):
        left = left0 + i * (box_w + gap)
        _add_rect(slide, left, top_s, box_w, Inches(5.7), fill=GREY_LIGHT)
        _add_rect(slide, left, top_s, box_w, Inches(0.9), fill=BLUE_DARK)
        _add_text_box(slide, left + Inches(0.12), top_s + Inches(0.08),
                      box_w - Inches(0.2), Inches(0.4),
                      pt_s, bold=True, size=13, color=WHITE)
        _add_text_box(slide, left + Inches(0.12), top_s + Inches(0.48),
                      box_w - Inches(0.2), Inches(0.35),
                      en_s, italic=True, size=10, color=BLUE_LIGHT)

        _add_text_box(slide, left + Inches(0.12), top_s + Inches(1.0),
                      box_w - Inches(0.25), Inches(1.3),
                      pt_d, size=11, color=DARK_TEXT)
        _divider(slide, top_s + Inches(2.35), color=GREY_MID)
        _add_text_box(slide, left + Inches(0.12), top_s + Inches(2.45),
                      box_w - Inches(0.25), Inches(1.3),
                      en_d, italic=True, size=10, color=BLUE_MID)

        _add_rect(slide, left + Inches(0.12), top_s + Inches(3.9),
                  box_w - Inches(0.25), Inches(1.5), fill=conf_color)
        _add_text_box(slide, left + Inches(0.12), top_s + Inches(4.25),
                      box_w - Inches(0.25), Inches(1.0),
                      conf, bold=True, size=16, color=WHITE,
                      align=PP_ALIGN.CENTER)

    _section_label(slide, "7  Validação | Validation")


def slide_roteiro(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "Roteiro de Melhoria — Dirigido por Lacunas de Confiança",
               "Improvement Roadmap — Driven by Confidence Gaps")

    phases = [
        ("Phase A", "Validação da cadeia\nde raciocínio",
         "Reasoning chain\nvalidation",
         [
             "A1 Fix BM25 + Rio Manual hints  ✓",
             "A2 Corpus sintético  ✓",
             "A3 Retrieval quality  ✓",
             "A4 Prompt quality  ✓",
             "A5 Slides + narrativa  ← aqui",
         ],
         "Jun 2026  |  CONCLUÍDA", GREEN_OK),
        ("Phase B", "Validação com\ndocumento real",
         "Real-document\nvalidation",
         [
             "B1 Obter processo PPP completo",
             "B2 Avaliar Ação 1 no real",
             "B3 Revisão do auditor",
             "B4 Comparação de modelos",
         ],
         "Jul 2026  |  Pendente", ORANGE),
        ("Phase C", "Backlog de\nmelhoria",
         "Improvement\nbacklog",
         [
             "Extensão: 45 ações restantes",
             "Validação frontend c/ usuários",
             "Benchmarks de performance",
             "Consistência entre avaliações",
         ],
         "Não agendado  |  Após Phase B", RED_LOW),
    ]
    box_w = Inches(3.9)
    gap = Inches(0.37)
    left0 = Inches(0.35)
    top = Inches(1.3)
    for i, (phase, pt, en, items, status, status_color) in enumerate(phases):
        left = left0 + i * (box_w + gap)
        _add_rect(slide, left, top, box_w, Inches(5.8), fill=GREY_LIGHT)
        _add_rect(slide, left, top, box_w, Inches(1.15), fill=BLUE_DARK)
        _add_text_box(slide, left + Inches(0.12), top + Inches(0.08),
                      box_w - Inches(0.2), Inches(0.4),
                      phase, bold=True, size=20, color=GOLD)
        _add_text_box(slide, left + Inches(0.12), top + Inches(0.48),
                      box_w - Inches(0.2), Inches(0.35),
                      pt, bold=True, size=13, color=WHITE)
        _add_text_box(slide, left + Inches(0.12), top + Inches(0.82),
                      box_w - Inches(0.2), Inches(0.28),
                      en, italic=True, size=10, color=BLUE_LIGHT)

        _bullet_block(slide, left + Inches(0.12), top + Inches(1.28),
                      box_w - Inches(0.2), Inches(3.2),
                      items, size=11)

        _add_rect(slide, left + Inches(0.12), top + Inches(4.6),
                  box_w - Inches(0.25), Inches(0.75), fill=status_color)
        _add_text_box(slide, left + Inches(0.12), top + Inches(4.72),
                      box_w - Inches(0.25), Inches(0.5),
                      status, bold=True, size=11, color=WHITE,
                      align=PP_ALIGN.CENTER)

    _section_label(slide, "8  Roteiro | Roadmap")


def slide_proximos_passos(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, fill=WHITE)
    _title_bar(slide,
               "Próximos Passos | Next Steps",
               "")

    # Single key question
    _add_rect(slide, Inches(0.35), Inches(1.3), Inches(12.6), Inches(1.2),
              fill=BLUE_DARK)
    _add_text_box(slide, Inches(0.55), Inches(1.4), Inches(12), Inches(0.5),
                  "Questão central: a cadeia de raciocínio é defensável com documentos reais?",
                  bold=True, size=20, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(0.55), Inches(1.88), Inches(12), Inches(0.38),
                  "Central question: is the reasoning chain defensible with real documents?",
                  italic=True, size=13, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    next_steps = [
        ("B1", "Obter processo PPP real completo\nObtain a complete real PPP process",
         "Acesso ao acervo municipal\nAccess to municipal archive",
         BLUE_MID),
        ("B2", "Avaliar Ação 1 no processo real\nEvaluate Action 1 on real process",
         "Comparar retrieval real vs. sintético\nCompare real vs. synthetic retrieval",
         BLUE_DARK),
        ("B3", "Revisão do auditor no resultado real\nAuditor review of real result",
         "Identificar padrões de substituição\nIdentify override patterns",
         BLUE_MID),
        ("B4", "Comparar Mistral vs. Groq\nCompare Mistral vs. Groq",
         "Avaliar consistência e qualidade\nAssess consistency and quality",
         BLUE_DARK),
    ]
    box_w = Inches(2.9)
    gap = Inches(0.3)
    left0 = Inches(0.35)
    top_n = Inches(2.65)
    for i, (num, action, dep, color) in enumerate(next_steps):
        left = left0 + i * (box_w + gap)
        _add_rect(slide, left, top_n, box_w, Inches(3.8), fill=GREY_LIGHT)
        _add_rect(slide, left, top_n, Inches(0.55), Inches(3.8), fill=color)
        _add_text_box(slide, left + Inches(0.08), top_n + Inches(1.5),
                      Inches(0.4), Inches(0.45),
                      num, bold=True, size=16, color=WHITE,
                      align=PP_ALIGN.CENTER)
        _add_text_box(slide, left + Inches(0.65), top_n + Inches(0.2),
                      box_w - Inches(0.8), Inches(1.4),
                      action, bold=True, size=12, color=DARK_TEXT)
        _divider(slide, top_n + Inches(1.7), color=GREY_MID)
        _add_text_box(slide, left + Inches(0.65), top_n + Inches(1.85),
                      box_w - Inches(0.8), Inches(1.7),
                      dep, size=10, italic=True, color=BLUE_MID)

    # Closing
    _add_rect(slide, Inches(0.35), Inches(6.6), Inches(12.6), Inches(0.65),
              fill=BLUE_DARK)
    _add_text_box(slide, Inches(0.55), Inches(6.68), Inches(12), Inches(0.45),
                  "Repositório: github.com/Sjacuru/Maturity_Check  ·  Contato: sjacuru@gmail.com",
                  size=12, color=WHITE, align=PP_ALIGN.CENTER)

    _section_label(slide, "8  Roteiro | Roadmap")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_title(prs)
    slide_agenda(prs)
    slide_problema(prs)
    slide_ipmp(prs)
    slide_solucao(prs)
    slide_fluxo(prs)
    slide_modulos(prs)
    slide_cascata(prs)
    slide_prompt(prs)
    slide_auditor(prs)
    slide_estado_atual(prs)
    slide_validacao(prs)
    slide_confianca(prs)
    slide_roteiro(prs)
    slide_proximos_passos(prs)

    out = Path(__file__).parent.parent / "docs" / "PPP_Maturity_Check_Slides.pptx"
    prs.save(str(out))
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
