# -*- coding: utf-8 -*-
"""
Create the 5-slide professor progress presentation for PPP Maturity Check.

Output: docs/PPP_Maturity_Check_Professor_5_Slides.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BLUE_DARK = RGBColor(0x1A, 0x3A, 0x5C)
BLUE_MID = RGBColor(0x2E, 0x6D, 0xA4)
BLUE_LIGHT = RGBColor(0xD0, 0xE4, 0xF7)
GOLD = RGBColor(0xC9, 0xA0, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1C, 0x1C, 0x1C)
MUTED = RGBColor(0x66, 0x75, 0x85)
GREY = RGBColor(0xF2, 0xF4, 0xF7)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xB9, 0x3A, 0x32)

W = Inches(13.33)
H = Inches(7.5)


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.7)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=16, color=TEXT,
             bold=False, italic=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=14, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.level = 0
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def title_bar(slide, number, title, subtitle):
    add_rect(slide, 0, 0, W, Inches(0.95), BLUE_DARK)
    add_rect(slide, 0, Inches(0.95), W, Inches(0.05), GOLD)
    add_text(slide, Inches(0.35), Inches(0.16), Inches(0.55), Inches(0.45),
             f"{number}", size=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.9), Inches(0.14), Inches(7.8), Inches(0.42),
             title, size=22, color=WHITE, bold=True)
    add_text(slide, Inches(0.92), Inches(0.55), Inches(11.7), Inches(0.25),
             subtitle, size=10.5, color=BLUE_LIGHT, italic=True)


def section_footer(slide, text):
    add_text(slide, Inches(9.5), Inches(7.1), Inches(3.45), Inches(0.25),
             text, size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


def metric(slide, left, label, value, detail, color=BLUE_MID):
    add_rect(slide, left, Inches(1.28), Inches(3.0), Inches(1.2), color)
    add_text(slide, left, Inches(1.34), Inches(3.0), Inches(0.45),
             value, size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.15), Inches(1.92), Inches(2.7), Inches(0.22),
             label, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.15), Inches(2.17), Inches(2.7), Inches(0.2),
             detail, size=8.5, color=BLUE_LIGHT, italic=True, align=PP_ALIGN.CENTER)


def pill(slide, left, top, width, text, color):
    add_rect(slide, left, top, width, Inches(0.42), color)
    add_text(slide, left + Inches(0.05), top + Inches(0.08), width - Inches(0.1),
             Inches(0.2), text, size=10, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, WHITE)
    title_bar(slide, "1", "Project Progress", "PPP Maturity Check: decision support for maturity assessment of PPP procurement documents")

    metric(slide, Inches(0.45), "Modules implemented", "7", "ingestion to frontend", BLUE_MID)
    metric(slide, Inches(3.65), "Automated tests", "300", "reported in project deck", GREEN)
    metric(slide, Inches(6.85), "IPMP scope", "1/46", "Action 1 completed", ORANGE)
    metric(slide, Inches(10.05), "Pipeline status", "E2E", "functional flow", BLUE_DARK)

    add_text(slide, Inches(0.55), Inches(2.95), Inches(5.8), Inches(0.35),
             "What has been built", size=15, color=BLUE_DARK, bold=True)
    add_bullets(slide, Inches(0.55), Inches(3.35), Inches(5.8), Inches(2.6), [
        "PDF extraction into provenance-preserving chunks.",
        "Retrieval cascade: document matching, BM25, regex, and vector fallback.",
        "LLM evaluation that proposes a 0 / 1 / 3 Maturity Score.",
        "FastAPI + Vue interface for upload, evidence review, and auditor decision."
    ], size=12.5)

    add_text(slide, Inches(6.9), Inches(2.95), Inches(5.8), Inches(0.35),
             "Current limitations", size=15, color=BLUE_DARK, bold=True)
    add_bullets(slide, Inches(6.9), Inches(3.35), Inches(5.8), Inches(2.6), [
        "Only IPMP Action 1 is fully encoded and evaluated.",
        "Validation so far is strongest on controlled or synthetic corpus.",
        "Real PPP document validation and model comparison remain next steps.",
        "The frontend is functional, but still needs user validation with auditors."
    ], size=12.5)
    section_footer(slide, "Progress")


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, WHITE)
    title_bar(slide, "2", "Refined Research Direction", "From broad automation to defensible, reproducible auditor support")

    add_rect(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(1.25), BLUE_DARK)
    add_text(slide, Inches(0.75), Inches(1.55), Inches(11.8), Inches(0.45),
             "Research question", size=13, color=GOLD, bold=True)
    add_text(slide, Inches(0.75), Inches(1.92), Inches(11.8), Inches(0.42),
             "Can an AI-assisted workflow produce traceable, reproducible, and auditor-verifiable PPP maturity assessments under the IPMP framework?",
             size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    labels = [
        ("Initial idea", "Automate maturity scoring across PPP documents.", MUTED),
        ("Refinement", "Prioritize one deeply validated IPMP Action before scaling to 46.", BLUE_MID),
        ("Contribution", "A forensic evidence package: retrieved chunks, prompt, reasoning, flags, and proposed score.", GREEN),
    ]
    top = Inches(3.0)
    for i, (head, body, color) in enumerate(labels):
        y = top + Inches(i * 1.15)
        pill(slide, Inches(0.7), y, Inches(2.0), head, color)
        add_text(slide, Inches(3.0), y - Inches(0.02), Inches(9.5), Inches(0.45),
                 body, size=15, color=TEXT)

    add_text(slide, Inches(0.7), Inches(6.35), Inches(11.8), Inches(0.38),
             "Next research step: validate Action 1 with a complete real PPP case, then compare model consistency and auditor override patterns.",
             size=14, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    section_footer(slide, "Research Direction")


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, WHITE)
    title_bar(slide, "3", "Methodology", "A staged pipeline with controlled validation before real-world assessment")

    steps = [
        ("1", "Ingest IPMP + Rio Manual", "Canonical JSON artifacts define the Acao, expected products, scoring rubric, and Rio-specific retrieval vocabulary."),
        ("2", "Extract document chunks", "PDF pages become chunks with filename, page, offset, OCR status, and source metadata."),
        ("3", "Retrieve evidence", "Deterministic lexical cascade runs first; vector search is fallback only when lexical retrieval returns no evidence."),
        ("4", "Evaluate with LLM", "Fixed prompt, temperature zero, sentinel output, parsed reasoning, proposed score, and uncertainty flags."),
        ("5", "Auditor validation", "Final score is accepted or overridden by the auditor with justification and evidence references.")
    ]
    left = Inches(0.45)
    top = Inches(1.35)
    for i, (num, head, body) in enumerate(steps):
        y = top + Inches(i * 1.05)
        add_rect(slide, left, y, Inches(0.62), Inches(0.62), BLUE_MID if i % 2 else BLUE_DARK)
        add_text(slide, left, y + Inches(0.11), Inches(0.62), Inches(0.3),
                 num, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, left + Inches(0.85), y - Inches(0.03), Inches(3.5), Inches(0.28),
                 head, size=13.5, color=BLUE_DARK, bold=True)
        add_text(slide, left + Inches(4.2), y - Inches(0.02), Inches(8.3), Inches(0.58),
                 body, size=11.7, color=TEXT)

    add_rect(slide, Inches(0.45), Inches(6.68), Inches(12.35), Inches(0.42), GREY)
    add_text(slide, Inches(0.65), Inches(6.77), Inches(12.0), Inches(0.18),
             "Validation plan: Phase A controlled corpus -> Phase B real PPP documents -> Phase C scale to more IPMP Actions.",
             size=11.5, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    section_footer(slide, "Methodology")


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, WHITE)
    title_bar(slide, "4", "Software Engineering Principles", "Design choices shaped by reproducibility, auditability, and maintainability")

    cards = [
        ("Reproducibility", "BM25, fixed prompts, and temperature-zero LLM calls make repeated evaluations comparable."),
        ("Traceability", "Every EvaluationResult stores evidence, provenance, prompts, raw response, reasoning, and status flags."),
        ("Separation of concerns", "Ingestion, extraction, retrieval, evaluation, assessment, API, and frontend have clear boundaries."),
        ("Human in the loop", "The system proposes; the auditor decides. Overrides require justification."),
        ("Progressive enrichment", "New Acoes can be added through source-of-truth artifacts rather than rewriting the core flow."),
        ("Testable contracts", "Pydantic models, protocol interfaces, and automated tests keep module behavior explicit.")
    ]
    for i, (head, body) in enumerate(cards):
        row, col = divmod(i, 3)
        x = Inches(0.45) + Inches(col * 4.25)
        y = Inches(1.35) + Inches(row * 2.45)
        add_rect(slide, x, y, Inches(3.9), Inches(1.9), GREY, line=BLUE_LIGHT)
        add_rect(slide, x, y, Inches(3.9), Inches(0.18), GOLD if i in (0, 1, 3) else BLUE_MID)
        add_text(slide, x + Inches(0.22), y + Inches(0.34), Inches(3.45), Inches(0.32),
                 head, size=14, color=BLUE_DARK, bold=True)
        add_text(slide, x + Inches(0.22), y + Inches(0.78), Inches(3.45), Inches(0.82),
                 body, size=11.4, color=TEXT)
    section_footer(slide, "Engineering")


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, WHITE)
    title_bar(slide, "5", "Ethical and Cultural Considerations", "Responsible use across workplace contexts, with a note on New Zealand and Māori culture")

    add_text(slide, Inches(0.55), Inches(1.25), Inches(5.9), Inches(0.34),
             "Ethical safeguards", size=15, color=BLUE_DARK, bold=True)
    add_bullets(slide, Inches(0.55), Inches(1.68), Inches(5.9), Inches(2.35), [
        "Maintain human accountability: final Maturity Score belongs to the auditor.",
        "Protect procurement data through access control, retention limits, and careful handling of confidential documents.",
        "Expose uncertainty, parse failures, and no-evidence states instead of hiding weak results.",
        "Avoid overclaiming: the system supports judgment; it does not replace professional or legal responsibility."
    ], size=12)

    add_text(slide, Inches(6.85), Inches(1.25), Inches(5.9), Inches(0.34),
             "Cultural fit for work environments", size=15, color=BLUE_DARK, bold=True)
    add_bullets(slide, Inches(6.85), Inches(1.68), Inches(5.9), Inches(2.35), [
        "Use inclusive language and transparent decision records so teams can challenge system outputs respectfully.",
        "Support local governance expectations, public-sector norms, and different documentation styles.",
        "Design training and rollout around shared learning, not surveillance of individual staff.",
        "Keep explanations understandable for technical and non-technical stakeholders."
    ], size=12)

    add_rect(slide, Inches(0.55), Inches(4.65), Inches(12.2), Inches(1.55), BLUE_DARK)
    add_text(slide, Inches(0.8), Inches(4.82), Inches(11.7), Inches(0.3),
             "New Zealand / Māori culture note", size=13, color=GOLD, bold=True)
    add_text(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.62),
             "In a New Zealand context, the project should respect Māori data interests and relationship-based ways of working. Practical adaptation would include early consultation with relevant Māori stakeholders, attention to Te Tiriti o Waitangi obligations, and values such as manaakitanga, whanaungatanga, and kaitiakitanga when public-service data or community impacts are involved.",
             size=12.2, color=WHITE)
    section_footer(slide, "Ethics & Culture")


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)

    out = Path(__file__).parent.parent / "docs" / "PPP_Maturity_Check_Professor_5_Slides.pptx"
    prs.save(out)
    print(f"Saved: {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
